import json
import os
from typing import Dict, List, Union, Any, Tuple
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import logging
from .tool_registry import ToolHandler
from path_manager import path_manager

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)





class PPTXGenerator:
    """
    PPTX生成器。
    """
    
    def __init__(self):
        """
        初始化生成器，创建空白演示文稿
        """
        self.prs = Presentation()
        # 16:9 宽屏
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        

    
    def _get_layout_settings(self):
        """获取布局设置"""
        return {
            # Layout Dimensions
            "margin_left": Inches(0.7),
            "margin_top": Inches(0.5),
            "margin_right": Inches(0.7),
            "margin_bottom": Inches(0.5),
            "column_gap": Inches(0.5),
            "header_height": Inches(0.7),
            "item_spacing": Inches(0.2), # 内容项之间的垂直间距
        }
    
    

    def create_cover_slide(self, title: str, subtitle: str = "") -> None:
        """创建封面幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)
        

        title_shape = slide.shapes.title
        title_shape.text = title
        p = title_shape.text_frame.paragraphs[0]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            p = subtitle_shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
    def create_section_slide(self, section_title: str, section_subtitle: str = "") -> None:
        """创建章节页幻灯片"""
        slide_layout = self.prs.slide_layouts[2] # 节标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide)

        title_shape = slide.shapes.title
        title_shape.text = section_title
        p = title_shape.text_frame.paragraphs[0]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        body_shape = slide.placeholders[1]
        if section_subtitle:
            body_shape.text = section_subtitle
            p = body_shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
        else:
            # 删除未使用的占位符
            sp = body_shape._element
            sp.getparent().remove(sp)

    def create_content_slide(self, slide_data: Dict[str, Any]) -> None:
        """根据指定的布局创建正文页幻灯片"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        

        layout_settings = self._get_layout_settings()
        
        # 添加幻灯片标题
        current_top = layout_settings["margin_top"]
        if "title" in slide_data:
            title_box = slide.shapes.add_textbox(
                layout_settings["margin_left"], current_top, 
                self.prs.slide_width - layout_settings["margin_left"] - layout_settings["margin_right"], 
                layout_settings["header_height"]
            )
            tf = title_box.text_frame
            tf.text = slide_data["title"]
            p = tf.paragraphs[0]
            p.font.bold = True
            
            current_top += layout_settings["header_height"] + layout_settings["item_spacing"]

        # 根据布局分发内容
        layout_type = slide_data.get("layout", "default")
        contents = slide_data.get("contents", [])
        
        if layout_type == "two_column":
            self._layout_two_column(slide, contents, current_top)
        else: # 默认为瀑布流布局
            self._layout_default(slide, contents, current_top)

    def _layout_default(self, slide, contents: List[Dict], start_top: float):
        """处理默认的瀑布流布局"""
        layout_settings = self._get_layout_settings()
        current_top = start_top
        available_width = self.prs.slide_width - layout_settings["margin_left"] - layout_settings["margin_right"]
        
        for item in contents:
            item_height = self._add_content_item(slide, item, 
                layout_settings["margin_left"], current_top, 
                available_width
            )
            current_top += item_height + layout_settings["item_spacing"]

    def _layout_two_column(self, slide, contents: List[Dict], start_top: float):
        """处理两栏布局"""
        layout_settings = self._get_layout_settings()
        col1_items = [c for c in contents if c.get("column") == 1]
        col2_items = [c for c in contents if c.get("column") == 2]

        col_width = (self.prs.slide_width - layout_settings["margin_left"] - layout_settings["margin_right"] - layout_settings["column_gap"]) / 2
        
        # 处理第一栏
        current_top_col1 = start_top
        col1_x = layout_settings["margin_left"]
        for item in col1_items:
            item_height = self._add_content_item(slide, item, col1_x, current_top_col1, col_width)
            current_top_col1 += item_height + layout_settings["item_spacing"]

        # 处理第二栏
        current_top_col2 = start_top
        col2_x = layout_settings["margin_left"] + col_width + layout_settings["column_gap"]
        for item in col2_items:
            item_height = self._add_content_item(slide, item, col2_x, current_top_col2, col_width)
            current_top_col2 += item_height + layout_settings["item_spacing"]

    def _add_content_item(self, slide, item: Dict[str, Any], x: float, y: float, width: float) -> float:
        """
        内容分发器：根据内容类型调用相应的添加方法
        返回: 添加的元素的高度
        """
        layout_settings = self._get_layout_settings()
        content_type = item.get("type")
        available_height = self.prs.slide_height - y - layout_settings["margin_bottom"]
        if available_height <= 0: return 0

        height = 0
        if content_type == "text":
            height = self._add_text_content(slide, item, x, y, width, available_height)
        elif content_type == "chart":
            height = self._add_chart_content(slide, item, x, y, width, available_height)
        elif content_type == "table":
            height = self._add_table_content(slide, item, x, y, width, available_height)
        # elif content_type == "image":  # 扩展点
        #     height = self._add_image_content(slide, item, x, y, width, available_height)
        else:
            logger.warning(f"未知的内容类型: {content_type}")
        
        return height

    def _add_text_content(self, slide, item: Dict, x: float, y: float, width: float, max_height: float) -> float:
        """在一个指定的矩形区域内添加文本内容"""
        # 预估高度，如果超出则截断
        estimated_lines = len(item.get("text", "").split('\n')) + len(item.get("bullet_points", []))
        estimated_height = Inches(estimated_lines * 0.3) # 粗略估计
        height = min(estimated_height, max_height, Inches(4)) # 最大高度为4英寸或可用高度

        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # 清除默认段落
        p = tf.paragraphs[0]
        p.text = ""

        if "text" in item:
            p = tf.add_paragraph()
            p.text = item["text"]
        
        if "bullet_points" in item:
            for i, point in enumerate(item["bullet_points"]):
                p = tf.add_paragraph()
                # 支持多级列表：如果point是字典，则提取text和level；否则作为简单字符串处理
                if isinstance(point, dict):
                    p.text = point.get("text", "")
                    p.level = point.get("level", 1)
                else:
                    p.text = str(point)
                    p.level = 1

        
        return height

    def _add_chart_content(self, slide, item: Dict, x: float, y: float, width: float, max_height: float) -> float:
        """在一个指定的矩形区域内添加图表"""
        height = min(max_height, Inches(4)) # 限制图表最大高度

        chart_data = CategoryChartData()
        chart_data.categories = item.get("data", {}).get("categories", [])
        for series_name, values in item.get("data", {}).get("series", {}).items():
            chart_data.add_series(series_name, values)
            
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED, 
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE, 
            "area": XL_CHART_TYPE.AREA_STACKED
        }
        xl_chart_type = chart_type_map.get(item.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)

        try:
            # 创建图表
            chart_shape = slide.shapes.add_chart(xl_chart_type, x, y, width, height, chart_data)
            chart = chart_shape.chart
            
            # 设置图表标题
            if item.get("chart_title"):
                chart.has_title = True
                chart.chart_title.text_frame.text = item["chart_title"]

            
            logger.info(f"图表创建成功: {item.get('chart_title', '未命名图表')}")
            
        except Exception as e:
            logger.error(f"创建图表失败: {e}")
            # 只有在图表创建失败时才显示错误信息
            textbox = slide.shapes.add_textbox(x, y, width, Inches(1))
            textbox.text_frame.text = f"图表生成失败: {item.get('chart_title', '')}"

            return Inches(1)
        
        return height

    def _add_table_content(self, slide, item: Dict, x: float, y: float, width: float, max_height: float) -> float:
        """在一个指定的矩形区域内添加表格"""
        headers = item.get("headers", [])
        rows = item.get("rows", [])
        table_title = item.get("table_title", "")
        
        if not headers or not rows:
            logger.warning("表格数据不完整，跳过表格创建")
            return Inches(0.5)
        
        # 计算表格尺寸
        cols = len(headers)
        table_rows = len(rows) + 1  # +1 for header row
        
        # 限制表格高度
        estimated_row_height = Inches(0.4)
        estimated_table_height = estimated_row_height * table_rows
        
        # 如果有标题，为标题预留空间
        title_height = Inches(0.3) if table_title else Inches(0)
        total_height = min(estimated_table_height + title_height, max_height)
        table_height = total_height - title_height
        
        current_y = y
        
        # 添加表格标题（如果有）
        if table_title:
            title_box = slide.shapes.add_textbox(x, current_y, width, title_height)
            tf = title_box.text_frame
            tf.text = table_title
            p = tf.paragraphs[0]
            p.font.bold = True
            current_y += title_height
        
        try:
            # 创建表格
            table_shape = slide.shapes.add_table(table_rows, cols, x, current_y, width, table_height)
            table = table_shape.table
            

            
            # 设置表头
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                # 设置表头样式
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER

            
            # 填充数据行（斑马纹效果）
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols:  # 确保不超出列数
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_data) if cell_data is not None else ""
                        # 设置数据行样式
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.LEFT

            
            logger.info(f"表格创建成功: {table_title or '未命名表格'} ({table_rows}行 x {cols}列)")
            
        except Exception as e:
            logger.error(f"创建表格失败: {e}")
            # 表格创建失败时显示错误信息
            error_box = slide.shapes.add_textbox(x, current_y, width, Inches(1))
            error_box.text_frame.text = f"表格生成失败: {table_title or '未命名表格'}"

            return title_height + Inches(1)
        
        return total_height

    def generate_from_json(self, json_content: Union[str, Dict]):
        """从JSON内容生成完整的演示文稿"""
        data = json.loads(json_content) if isinstance(json_content, str) else json_content
        
        if "slides" not in data:
            raise ValueError("JSON必须包含'slides'字段")
            
        for i, slide_data in enumerate(data["slides"]):
            slide_type = slide_data.get("type", "content")
            logger.info(f"正在生成幻灯片 {i+1}: 类型 '{slide_type}'")
            try:
                if slide_type == "cover":
                    self.create_cover_slide(slide_data.get("title", ""), slide_data.get("subtitle", ""))
                elif slide_type == "section":
                    self.create_section_slide(slide_data.get("title", ""), slide_data.get("subtitle", ""))
                elif slide_type == "content":
                    self.create_content_slide(slide_data)
                else:
                    logger.warning(f"跳过未知的幻灯片类型: {slide_type}")
            except Exception as e:
                logger.error(f"生成第 {i+1} 张幻灯片时出错: {e}", exc_info=True)
                # 可选择在这里停止或继续
                # raise

    def save(self, filename: str) -> str:
        """保存PPTX文件"""
        # 使用path_manager来处理路径
        output_path = Path(filename)
        if not output_path.is_absolute():
             output_path = path_manager.get_output_path(filename)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info(f"PPTX文件已成功保存至: {output_path}")
        return str(output_path)

def create_pptx_from_json(json_content: Union[str, Dict], output_filename: str) -> str:
    """
    主函数：从JSON生成PPTX文件。
    
    参数:
        json_content: JSON内容 (字符串或字典)。
        output_filename: 输出文件名。
        
    返回:
        生成的文件路径。
    """
    try:
        generator = PPTXGenerator()
        generator.generate_from_json(json_content)
        return generator.save(output_filename)
    except Exception as e:
        logger.error(f"生成PPTX时发生严重错误: {e}")
        raise


class PptCreationTool(ToolHandler):
    """PPT创建工具"""
    
    @property
    def name(self) -> str:
        return "create_pptx_presentation"
    
    def execute(self, args: Dict[str, Any], context: Dict[str, Any]) -> Tuple[bool, str]:
        """
        创建PPT文件
        
        参数:
            args: 包含json_content和output_filename的参数字典
            context: 执行上下文（当前未使用）
            
        返回:
            (成功标志, 执行结果)
        """
        json_content = args.get("json_content", {})
        output_filename = args.get("output_filename", "output")
        
        if not json_content:
            return False, json.dumps({
                "error": "缺少PPT内容数据"
            }, ensure_ascii=False)
        
        try:
            file_path = create_pptx_from_json(json_content, output_filename)
            return True, json.dumps({
                "success": True,
                "file_path": file_path,
                "message": f"PPT文件已成功生成：{file_path}"
            }, ensure_ascii=False)
        except Exception as e:
            return False, json.dumps({
                "success": False,
                "error": str(e)
            }, ensure_ascii=False)



